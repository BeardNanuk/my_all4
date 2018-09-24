
# ze plot functions
# create on Fri Sep 21 20:57:01 UTC 2018

import argparse, os, sys
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec

# textwrap for plotting text
import textwrap

# generating PPT slides
from pptx import Presentation
from pptx.util import Inches


from seisflows.config import config, loadpy, tilde_expand, Dict
from seisflows.tools import unix

import datetime

import glob
import os

from scipy import ndimage


def textplot(struct_variable,filename=None,flag_close=None,flag_save_fig=None):
    
    if filename is None:
	filename = 0

    if flag_close is None:
	flag_close = 1 

    if flag_save_fig is None:
	flag_save_fig = 1

    attributes = [a for a in dir(struct_variable)
                  if not (a.startswith('__') and a.endswith('__'))
                  and not (a == 'kernel')]

    fig = plt.figure(figsize=(10, 7)) 
    #fig.suptitle('bold figure suptitle', fontsize=14, fontweight='bold')

    ax = fig.add_subplot(111)
    fig.subplots_adjust(top=0.85)
    ax.set_title(str('parameters'))

    ax.set_xlabel('xlabel')
    ax.set_ylabel('ylabel')

    spc = 0.26

    for i in range(len(attributes)):
        offset = spc * (i) 
        str_temp = attributes[i] + ' = ' + str((getattr(struct_variable, attributes[i])))
        #str_temp = attributes[i] + ' = ' + str("%.5f" % (getattr(struct_variable, attributes[i])))
       ###ax.text(0.3, 9.7-offset,textwrap.shorten(str_temp, 110), fontsize=11)
        ax.text(0.3, 9.7-offset, str_temp, fontsize=9)

    ax.axis([0, 10, 0, 10])
    #plt.axis('off')
    if flag_save_fig==1:
        plt.savefig(filename,format='png', dpi=500)
    #plt.show()
    if flag_close == 1:
        plt.close()



def add_slide_ze(img_path,filename_pptx,left_start=None,top_start=None,width=None,height=None):

    if left_start is None:
       left_start = 0

    if top_start is None:
       top_start = 0.9  

    if width is None:
       width = 9

    if height is None:
       height = 5.7

    prs = Presentation(filename_pptx)

    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Inches(left_start)
    top = Inches(top_start)
    width = Inches(width)
    height = Inches(height)


    pic = slide.shapes.add_picture(img_path, left, top,width=width,height=height)

    # left = Inches(5)
    # height = Inches(5.5)
    # pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save(filename_pptx)




